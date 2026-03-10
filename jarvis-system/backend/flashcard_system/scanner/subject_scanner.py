"""
Subject Scanner and Material Processor for Adaptive Learning.
Handles directory scanning for subject materials and text extraction.
"""

import os
import logging
from pathlib import Path
from typing import List, Dict, Any, Optional
import PyPDF2
import docx
from datetime import datetime

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class MaterialProcessor:
    """Extracts text content from various educational materials (PDF, DOCX, TXT, etc.)."""
    
    SUPPORTED_EXTENSIONS = {'.pdf', '.docx', '.txt', '.md', '.ppt', '.pptx'}

    @staticmethod
    def extract_text(file_path: Path) -> str:
        """
        Extract text from a file based on its extension.
        
        Args:
            file_path: Path to the material file
            
        Returns:
            Extracted text content
        """
        suffix = file_path.suffix.lower()
        try:
            if suffix == '.pdf':
                return MaterialProcessor._extract_pdf(file_path)
            elif suffix == '.docx':
                return MaterialProcessor._extract_docx(file_path)
            elif suffix in ['.txt', '.md']:
                return MaterialProcessor._extract_text(file_path)
            elif suffix in ['.ppt', '.pptx']:
                # For now, return a placeholder or attempt basic extraction if library exists
                return MaterialProcessor._extract_pptx(file_path)
            else:
                logger.warning(f"Unsupported file type: {suffix}")
                return ""
        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {e}")
            return ""

    @staticmethod
    def _extract_pdf(file_path: Path) -> str:
        """Extract text from PDF using PyPDF2."""
        text = ""
        with open(file_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                extracted = page.extract_text()
                if extracted:
                    text += extracted + "\n"
        return text

    @staticmethod
    def _extract_docx(file_path: Path) -> str:
        """Extract text from DOCX using python-docx."""
        doc = docx.Document(file_path)
        return "\n".join([p.text for p in doc.paragraphs])

    @staticmethod
    def _extract_text(file_path: Path) -> str:
        """Extract text from plain text files."""
        return file_path.read_text(encoding='utf-8', errors='ignore')

    @staticmethod
    def _extract_pptx(file_path: Path) -> str:
        """Extract text from PPT/PPTX. (Requires python-pptx)."""
        try:
            import pptx
            prs = pptx.Presentation(file_path)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            return "\n".join(text_runs)
        except ImportError:
            logger.warning("python-pptx not installed, skipping PPT extraction.")
            return f"[PPT Content Placeholder for {file_path.name}]"
        except Exception as e:
            logger.error(f"Error extracting PPTX: {e}")
            return ""

class SubjectScanner:
    """Scans root directory for subject folders and their educational materials."""

    def __init__(self, root_dir: Optional[Path] = None):
        """
        Initialize scanner.
        
        Args:
            root_dir: Root directory to scan
        """
        # Set default to project root/data/subjects
        project_root = Path(__file__).parent.parent.parent
        self.root_dir = root_dir or (project_root / "data" / "subjects")
        
        if not self.root_dir.exists():
            try:
                self.root_dir.mkdir(parents=True, exist_ok=True)
                logger.info(f"Created material root directory at {self.root_dir}")
            except Exception as e:
                logger.error(f"Could not create root directory {self.root_dir}: {e}")

    def scan_all_subjects(self) -> Dict[str, List[Path]]:
        """
        Scan the root directory for all subjects and their files.
        
        Returns:
            Dictionary mapping subject names to list of file paths.
        """
        subjects = {}
        if not self.root_dir.exists():
            return {}

        for item in self.root_dir.iterdir():
            if item.is_dir():
                subject_name = item.name
                # Find all supported files in this subject folder (recursive)
                files = []
                for file_path in item.rglob('*'):
                    if file_path.is_file() and file_path.suffix.lower() in MaterialProcessor.SUPPORTED_EXTENSIONS:
                        files.append(file_path)
                
                if files:
                    subjects[subject_name] = files
        
        return subjects

    def get_subject_content(self, subject_name: str) -> str:
        """
        Get all text content for a specific subject by processing its files.
        
        Args:
            subject_name: Name of the subject (folder name)
            
        Returns:
            Aggregated text content
        """
        subject_dir = self.root_dir / subject_name
        logger.info(f"SubjectScanner probing directory: {subject_dir.absolute()}")
        try:
            with open("/tmp/jarvis_debug.txt", "w") as f:
                f.write(f"Probing: {subject_dir.absolute()}\nExists: {subject_dir.exists()}\nIs_dir: {subject_dir.is_dir()}\n")
        except:
            pass
        if not subject_dir.exists() or not subject_dir.is_dir():
            logger.warning(f"Subject directory {subject_name} not found. Path: {subject_dir.absolute()}")
            return ""

        all_text = []
        for file_path in subject_dir.rglob('*'):
            if file_path.is_file() and file_path.suffix.lower() in MaterialProcessor.SUPPORTED_EXTENSIONS:
                logger.info(f"Processing material: {file_path}")
                text = MaterialProcessor.extract_text(file_path)
                if text.strip():
                    all_text.append(f"--- SOURCE: {file_path.name} ---\n{text}")
        
        return "\n\n".join(all_text)

    def list_subjects(self) -> List[Dict[str, Any]]:
        """List all subjects with metadata."""
        subjects_data = []
        scanned = self.scan_all_subjects()
        
        for name, files in scanned.items():
            subjects_data.append({
                "name": name,
                "file_count": len(files),
                "last_modified": datetime.fromtimestamp((self.root_dir / name).stat().st_mtime).isoformat(),
                "path": str(self.root_dir / name)
            })
            
        return subjects_data

if __name__ == "__main__":
    # Test scanner
    scanner = SubjectScanner()
    subjects = scanner.list_subjects()
    print(f"Scanned Subjects: {subjects}")
