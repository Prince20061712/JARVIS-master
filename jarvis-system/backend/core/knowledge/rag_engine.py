import os
try:
    import chromadb
except ImportError:
    chromadb = None
    print("⚠️ ChromaDB module not found. RAG features will be disabled.")

from sentence_transformers import SentenceTransformer
import fitz  # PyMuPDF
from pptx import Presentation
from PIL import Image
import pytesseract

class EngineeringRAGEngine:
    def __init__(self, knowledge_base_path="engineering_knowledge"):
        self.kb_path = knowledge_base_path
        self.chroma_client = None
        self.collection = None
        
        if chromadb:
            try:
                self.chroma_client = chromadb.PersistentClient(path="chroma_db")
                self.collection = self.chroma_client.get_or_create_collection(name="engineering_knowledge")
            except Exception as e:
                print(f"⚠️ Failed to initialize ChromaDB client: {e}")
        
        # Load lightweight embedding model
        try:
            self.embedding_model = SentenceTransformer('all-MiniLM-L6-v2') 
            print("🔧 Engineering RAG Engine Initialized")
        except Exception as e:
             self.embedding_model = None
             print(f"⚠️ Failed to load embedding model: {e}")

    def ingest_document(self, file_path, subject):
        """Ingest a document (PDF, PPTX, image) into the vector store"""
        if not os.path.exists(file_path):
            return f"File not found: {file_path}"
            
        ext = os.path.splitext(file_path)[1].lower()
        chunks = []
        metadatas = []
        ids = []
        
        try:
            if ext == '.pdf':
                doc = fitz.open(file_path)
                for page_num, page in enumerate(doc):
                    text = page.get_text()
                    self._process_text_to_chunks(text, chunks, metadatas, ids, file_path, subject, page_num + 1)
            elif ext in ['.pptx']:
                prs = Presentation(file_path)
                for slide_num, slide in enumerate(prs.slides):
                    text = ""
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                    self._process_text_to_chunks(text, chunks, metadatas, ids, file_path, subject, slide_num + 1)
            elif ext in ['.jpg', '.jpeg', '.png']:
                image = Image.open(file_path)
                text = pytesseract.image_to_string(image)
                self._process_text_to_chunks(text, chunks, metadatas, ids, file_path, subject, 1)
            else:
                return f"Unsupported file type: {ext}"
        except Exception as e:
            return f"Error processing file: {e}"

        if chunks and self.collection and self.embedding_model:
            embeddings = self.embedding_model.encode(chunks).tolist()
            self.collection.add(
                documents=chunks,
                embeddings=embeddings,
                metadatas=metadatas,
                ids=ids
            )
            return f"✅ Ingested {len(chunks)} chunks from {os.path.basename(file_path)}"
        return "⚠️ No valid text found in document."

    def _process_text_to_chunks(self, text, chunks, metadatas, ids, file_path, subject, page_or_slide_num):
        """Helper to break text into chunks and append to lists"""
        if not text.strip():
            return
            
        page_chunks = [text[i:i+1000] for i in range(0, len(text), 800)] # 20% overlap
        
        for i, chunk in enumerate(page_chunks):
            if len(chunk.strip()) < 50: continue # Skip empty chunks
            
            chunks.append(chunk)
            metadatas.append({
                "source": os.path.basename(file_path),
                "page": page_or_slide_num,
                "subject": subject
            })
            ids.append(f"{os.path.basename(file_path)}_p{page_or_slide_num}_{i}")

    def retrieve_context(self, query, subject=None, n_results=4):
        """Retrieve relevant context for a query"""
        if not self.collection or not self.embedding_model:
            return ""

        query_embedding = self.embedding_model.encode([query]).tolist()
        
        where_filter = {"subject": subject} if subject else None
        
        results = self.collection.query(
            query_embeddings=query_embedding,
            n_results=n_results,
            where=where_filter
        )
        
        context_blocks = []
        if results['documents']:
            for i, doc in enumerate(results['documents'][0]):
                meta = results['metadatas'][0][i]
                source = f"[Source: {meta['source']}, Page {meta['page']}]"
                context_blocks.append(f"{source}\n{doc}")
                
        return "\n\n".join(context_blocks)

if __name__ == "__main__":
    # Test
    rag = EngineeringRAGEngine()
    print("RAG Engine Ready.")
