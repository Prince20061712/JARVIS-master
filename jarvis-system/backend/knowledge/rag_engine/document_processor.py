"""
Advanced document processor with multi-format support and intelligent preprocessing
"""

import os
import re
import hashlib
import asyncio
from typing import List, Dict, Any, Optional, Union, BinaryIO
from datetime import datetime
from pathlib import Path
import logging
from concurrent.futures import ThreadPoolExecutor
import json

# Document parsing libraries
import PyPDF2
from docx import Document
import markdown
from bs4 import BeautifulSoup
import pandas as pd
import yaml
import csv
from io import StringIO

# Text preprocessing
import nltk
from nltk.corpus import stopwords
from nltk.stem import WordNetLemmatizer
import re

# Try loading spaCy model
try:
    import spacy
    nlp = spacy.load("en_core_web_sm")
except:
    nlp = None

from .chunking_strategies import ChunkingStrategies, Chunk, ChunkingStrategy

# Download NLTK data
try:
    nltk.data.find('corpora/stopwords')
    nltk.data.find('corpora/wordnet')
except LookupError:
    nltk.download('stopwords')
    nltk.download('wordnet')

class DocumentType:
    """Supported document types"""
    PDF = "pdf"
    DOCX = "docx"
    TXT = "txt"
    MD = "md"
    HTML = "html"
    CSV = "csv"
    JSON = "json"
    YAML = "yaml"
    PPTX = "pptx"
    IMAGE = "image"
    UNKNOWN = "unknown"

class DocumentProcessor:
    """
    Advanced document processor with support for multiple file formats,
    intelligent preprocessing, and metadata extraction
    """
    
    def __init__(self, config: Optional[Dict[str, Any]] = None):
        """
        Initialize document processor with configuration
        
        Args:
            config: Configuration dictionary with processing parameters
        """
        self.config = config or {}
        self.chunking_strategies = ChunkingStrategies(config.get('chunking', {}))
        
        # Initialize NLP tools
        self.lemmatizer = WordNetLemmatizer()
        self.stop_words = set(stopwords.words('english'))
        
        # Load custom patterns for academic content
        self.academic_patterns = self._load_academic_patterns()
        
        # Thread pool for parallel processing
        self.executor = ThreadPoolExecutor(max_workers=self.config.get('max_workers', 4))
        
        self.logger = logging.getLogger(__name__)
        
    def process_document(
        self,
        file_path: Union[str, Path],
        chunking_strategy: ChunkingStrategy = ChunkingStrategy.SEMANTIC,
        metadata: Optional[Dict[str, Any]] = None,
        extract_metadata: bool = True,
        **kwargs
    ) -> Dict[str, Any]:
        """
        Process a document from file path
        
        Args:
            file_path: Path to the document
            chunking_strategy: Chunking strategy to use
            metadata: Additional metadata
            extract_metadata: Whether to extract metadata automatically
            **kwargs: Additional arguments
            
        Returns:
            Dictionary with processed document data
        """
        file_path = Path(file_path)
        if not file_path.exists():
            raise FileNotFoundError(f"Document not found: {file_path}")
        
        # Detect document type
        doc_type = self._detect_document_type(file_path)
        
        # Read and parse document
        raw_text, doc_metadata = self._parse_document(file_path, doc_type)
        
        # Extract metadata if requested
        if extract_metadata:
            doc_metadata.update(self._extract_metadata(raw_text, doc_type))
        
        # Merge with provided metadata
        if metadata:
            doc_metadata.update(metadata)
        
        # Preprocess text
        processed_text = self.preprocess_text(raw_text, doc_type)
        
        # Create chunks
        chunks = self.chunking_strategies.chunk_document(
            processed_text,
            strategy=chunking_strategy,
            metadata={
                'document_path': str(file_path),
                'document_type': doc_type,
                'filename': file_path.name,
                **doc_metadata
            },
            **kwargs
        )
        
        # Generate document ID
        doc_id = self._generate_document_id(file_path, raw_text)
        
        return {
            'document_id': doc_id,
            'file_path': str(file_path),
            'filename': file_path.name,
            'document_type': doc_type,
            'metadata': doc_metadata,
            'raw_text': raw_text,
            'processed_text': processed_text,
            'chunks': chunks,
            'num_chunks': len(chunks),
            'total_length': len(raw_text),
            'processed_at': datetime.now().isoformat()
        }
    
    def process_text(
        self,
        text: str,
        chunking_strategy: ChunkingStrategy = ChunkingStrategy.SEMANTIC,
        metadata: Optional[Dict[str, Any]] = None,
        **kwargs
    ) -> Dict[str, Any]:
        """
        Process raw text directly
        
        Args:
            text: Raw text to process
            chunking_strategy: Chunking strategy to use
            metadata: Additional metadata
            **kwargs: Additional arguments
            
        Returns:
            Dictionary with processed text data
        """
        # Extract metadata from text
        doc_metadata = self._extract_metadata(text, DocumentType.TXT)
        if metadata:
            doc_metadata.update(metadata)
        
        # Preprocess text
        processed_text = self.preprocess_text(text, DocumentType.TXT)
        
        # Create chunks
        chunks = self.chunking_strategies.chunk_document(
            processed_text,
            strategy=chunking_strategy,
            metadata=doc_metadata,
            **kwargs
        )
        
        # Generate document ID
        doc_id = hashlib.md5(text[:1000].encode()).hexdigest()[:16]
        
        return {
            'document_id': doc_id,
            'document_type': DocumentType.TXT,
            'metadata': doc_metadata,
            'raw_text': text,
            'processed_text': processed_text,
            'chunks': chunks,
            'num_chunks': len(chunks),
            'total_length': len(text),
            'processed_at': datetime.now().isoformat()
        }
    
    async def process_documents_parallel(
        self,
        file_paths: List[Union[str, Path]],
        chunking_strategy: ChunkingStrategy = ChunkingStrategy.SEMANTIC,
        **kwargs
    ) -> List[Dict[str, Any]]:
        """
        Process multiple documents in parallel
        
        Args:
            file_paths: List of file paths
            chunking_strategy: Chunking strategy to use
            **kwargs: Additional arguments
            
        Returns:
            List of processed document dictionaries
        """
        tasks = []
        for file_path in file_paths:
            task = asyncio.create_task(
                asyncio.to_thread(
                    self.process_document,
                    file_path,
                    chunking_strategy,
                    **kwargs
                )
            )
            tasks.append(task)
        
        results = await asyncio.gather(*tasks, return_exceptions=True)
        
        # Filter out exceptions
        processed = []
        for result in results:
            if isinstance(result, Exception):
                self.logger.error(f"Error processing document: {result}")
            else:
                processed.append(result)
        
        return processed
    
    def preprocess_text(
        self,
        text: str,
        doc_type: str = DocumentType.TXT,
        **kwargs
    ) -> str:
        """
        Advanced text preprocessing with multiple cleaning steps
        
        Args:
            text: Raw text to preprocess
            doc_type: Document type for specific preprocessing
            **kwargs: Additional preprocessing options
            
        Returns:
            Preprocessed text
        """
        if not text:
            return ""
        
        # Apply document-specific preprocessing
        if doc_type == DocumentType.HTML:
            text = self._clean_html(text)
        elif doc_type == DocumentType.MD:
            text = self._clean_markdown(text)
        
        # Basic cleaning
        text = self._basic_cleaning(text)
        
        # Remove extra whitespace
        text = re.sub(r'\s+', ' ', text)
        
        # Fix common OCR errors and formatting issues
        text = self._fix_formatting_errors(text)
        
        # Handle academic content specifically
        if kwargs.get('is_academic', True):
            text = self._enhance_academic_text(text)
        
        # Optional: Remove stopwords (usually not for RAG)
        if kwargs.get('remove_stopwords', False):
            text = self._remove_stopwords(text)
        
        # Optional: Lemmatization (usually not for RAG)
        if kwargs.get('lemmatize', False):
            text = self._lemmatize_text(text)
        
        return text.strip()
    
    def extract_entities(self, text: str) -> Dict[str, List[str]]:
        """
        Extract named entities from text using spaCy
        
        Args:
            text: Text to extract entities from
            
        Returns:
            Dictionary of entity types and values
        """
        if nlp is None:
            return {}
        
        doc = nlp(text[:100000])  # Limit to first 100k chars for performance
        
        entities = {
            'PERSON': [],
            'ORG': [],
            'GPE': [],
            'DATE': [],
            'MONEY': [],
            'PERCENT': [],
            'LAW': [],
            'PRODUCT': [],
            'EVENT': [],
            'WORK_OF_ART': []
        }
        
        for ent in doc.ents:
            if ent.label_ in entities:
                entities[ent.label_].append(ent.text)
        
        # Remove duplicates while preserving order
        for key in entities:
            entities[key] = list(dict.fromkeys(entities[key]))
        
        return entities
    
    def extract_key_phrases(self, text: str, top_n: int = 10) -> List[str]:
        """
        Extract key phrases using RAKE or similar algorithm
        
        Args:
            text: Text to extract phrases from
            top_n: Number of top phrases to return
            
        Returns:
            List of key phrases
        """
        # Simple implementation - in production use proper key phrase extraction
        words = text.lower().split()
        word_freq = {}
        
        # Calculate word frequencies
        for word in words:
            if word not in self.stop_words and len(word) > 3:
                word_freq[word] = word_freq.get(word, 0) + 1
        
        # Get top words by frequency
        top_words = sorted(word_freq.items(), key=lambda x: x[1], reverse=True)[:top_n]
        
        return [word for word, _ in top_words]
    
    def _detect_document_type(self, file_path: Path) -> str:
        """Detect document type from file extension"""
        extension = file_path.suffix.lower()
        
        type_map = {
            '.pdf': DocumentType.PDF,
            '.docx': DocumentType.DOCX,
            '.doc': DocumentType.DOCX,
            '.txt': DocumentType.TXT,
            '.md': DocumentType.MD,
            '.markdown': DocumentType.MD,
            '.html': DocumentType.HTML,
            '.htm': DocumentType.HTML,
            '.csv': DocumentType.CSV,
            '.json': DocumentType.JSON,
            '.yaml': DocumentType.YAML,
            '.yml': DocumentType.YAML,
            '.pptx': DocumentType.PPTX,
            '.jpg': DocumentType.IMAGE,
            '.jpeg': DocumentType.IMAGE,
            '.png': DocumentType.IMAGE
        }
        
        return type_map.get(extension, DocumentType.UNKNOWN)
    
    def _parse_document(self, file_path: Path, doc_type: str) -> tuple:
        """
        Parse document based on its type
        
        Returns:
            Tuple of (raw_text, metadata)
        """
        parsers = {
            DocumentType.PDF: self._parse_pdf,
            DocumentType.DOCX: self._parse_docx,
            DocumentType.TXT: self._parse_text,
            DocumentType.MD: self._parse_markdown,
            DocumentType.HTML: self._parse_html,
            DocumentType.CSV: self._parse_csv,
            DocumentType.JSON: self._parse_json,
            DocumentType.YAML: self._parse_yaml,
            DocumentType.PPTX: self._parse_pptx,
            DocumentType.IMAGE: self._parse_image
        }
        
        parser = parsers.get(doc_type, self._parse_text)
        return parser(file_path)
    
    def _parse_pdf(self, file_path: Path) -> tuple:
        """Parse PDF document"""
        text = ""
        metadata = {}
        
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                # Extract metadata
                if pdf_reader.metadata:
                    for key, value in pdf_reader.metadata.items():
                        metadata[key[1:]] = value  # Remove leading '/'
                
                # Extract text from each page
                for page_num, page in enumerate(pdf_reader.pages):
                    page_text = page.extract_text()
                    if page_text:
                        text += f"\n--- Page {page_num + 1} ---\n{page_text}"
                
        except Exception as e:
            self.logger.error(f"Error parsing PDF {file_path}: {e}")
        
        return text, {**metadata, 'num_pages': len(pdf_reader.pages)}
    
    def _parse_docx(self, file_path: Path) -> tuple:
        """Parse DOCX document"""
        text = ""
        metadata = {}
        
        try:
            doc = Document(file_path)
            
            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text:
                    text += para.text + "\n"
            
            # Extract tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = ' | '.join(cell.text for cell in row.cells)
                    text += row_text + "\n"
                text += "\n"
            
            # Get document properties
            if doc.core_properties:
                props = doc.core_properties
                metadata = {
                    'author': props.author,
                    'created': props.created,
                    'modified': props.modified,
                    'title': props.title,
                    'subject': props.subject
                }
                
        except Exception as e:
            self.logger.error(f"Error parsing DOCX {file_path}: {e}")
        
        return text, metadata
    
    def _parse_text(self, file_path: Path) -> tuple:
        """Parse plain text file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                text = file.read()
        except UnicodeDecodeError:
            # Try different encoding
            with open(file_path, 'r', encoding='latin-1') as file:
                text = file.read()
        
        return text, {}
    
    def _parse_markdown(self, file_path: Path) -> tuple:
        """Parse Markdown file"""
        with open(file_path, 'r', encoding='utf-8') as file:
            md_text = file.read()
        
        # Convert to HTML then extract text
        html = markdown.markdown(md_text)
        text = self._clean_html(html)
        
        return text, {'original_format': 'markdown'}
    
    def _parse_html(self, file_path: Path) -> tuple:
        """Parse HTML file"""
        with open(file_path, 'r', encoding='utf-8') as file:
            html_content = file.read()
        
        text = self._clean_html(html_content)
        
        return text, {'original_format': 'html'}
    
    def _parse_csv(self, file_path: Path) -> tuple:
        """Parse CSV file"""
        text = ""
        metadata = {'rows': 0, 'columns': 0}
        
        try:
            df = pd.read_csv(file_path)
            metadata['rows'] = len(df)
            metadata['columns'] = len(df.columns)
            metadata['column_names'] = list(df.columns)
            
            # Convert to text representation
            text = f"CSV Data: {metadata['rows']} rows, {metadata['columns']} columns\n"
            text += f"Columns: {', '.join(df.columns)}\n\n"
            text += df.to_string()
            
        except Exception as e:
            self.logger.error(f"Error parsing CSV {file_path}: {e}")
        
        return text, metadata
    
    def _parse_json(self, file_path: Path) -> tuple:
        """Parse JSON file"""
        with open(file_path, 'r', encoding='utf-8') as file:
            data = json.load(file)
        
        text = json.dumps(data, indent=2)
        
        return text, {'original_format': 'json'}
    
    def _parse_yaml(self, file_path: Path) -> tuple:
        """Parse YAML file"""
        with open(file_path, 'r', encoding='utf-8') as file:
            data = yaml.safe_load(file)
        
        text = yaml.dump(data, default_flow_style=False)
        
        return text, {'original_format': 'yaml'}
    
    def _parse_pptx(self, file_path: Path) -> tuple:
        """Parse PPTX document"""
        try:
            from pptx import Presentation
        except ImportError:
            self.logger.error("python-pptx not installed")
            return "", {}
            
        text = ""
        metadata = {}
        try:
            prs = Presentation(file_path)
            metadata['num_slides'] = len(prs.slides)
            for slide_num, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                text += "\n"
        except Exception as e:
            self.logger.error(f"Error parsing PPTX {file_path}: {e}")
        
        return text, metadata

    def _parse_image(self, file_path: Path) -> tuple:
        """Parse Image document"""
        try:
            from PIL import Image
            import pytesseract
        except ImportError:
            self.logger.error("Pillow or pytesseract not installed")
            return "", {}
            
        text = ""
        metadata = {}
        try:
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image)
        except Exception as e:
            self.logger.error(f"Error parsing Image {file_path}: {e}")
        
        return text, metadata
    
    def _clean_html(self, html_content: str) -> str:
        """Extract clean text from HTML"""
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()
        
        # Get text
        text = soup.get_text()
        
        # Clean up whitespace
        lines = (line.strip() for line in text.splitlines())
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        text = ' '.join(chunk for chunk in chunks if chunk)
        
        return text
    
    def _clean_markdown(self, md_text: str) -> str:
        """Clean markdown formatting"""
        # Remove markdown links: [text](url)
        text = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', md_text)
        
        # Remove images: ![alt](url)
        text = re.sub(r'!\[[^\]]*\]\([^\)]+\)', '', text)
        
        # Remove headers: # Header
        text = re.sub(r'^#+\s+', '', text, flags=re.MULTILINE)
        
        # Remove emphasis: *text* or **text**
        text = re.sub(r'(\*\*|__)(.*?)\1', r'\2', text)
        text = re.sub(r'(\*|_)(.*?)\1', r'\2', text)
        
        # Remove code blocks
        text = re.sub(r'```.*?```', '', text, flags=re.DOTALL)
        text = re.sub(r'`([^`]+)`', r'\1', text)
        
        return text
    
    def _basic_cleaning(self, text: str) -> str:
        """Basic text cleaning operations"""
        # Replace common problematic characters
        replacements = {
            '\u2018': "'", '\u2019': "'",  # Smart quotes
            '\u201c': '"', '\u201d': '"',
            '\u2013': '-', '\u2014': '--',
            '\u2026': '...',
            '\xa0': ' '  # Non-breaking space
        }
        
        for old, new in replacements.items():
            text = text.replace(old, new)
        
        # Remove control characters
        text = ''.join(char for char in text if ord(char) >= 32 or char in '\n\r\t')
        
        # Fix line breaks
        text = re.sub(r'\r\n', '\n', text)
        text = re.sub(r'\r', '\n', text)
        
        # Remove multiple line breaks
        text = re.sub(r'\n{3,}', '\n\n', text)
        
        return text
    
    def _fix_formatting_errors(self, text: str) -> str:
        """Fix common OCR and formatting errors"""
        # Fix hyphenated line breaks
        text = re.sub(r'(\w+)-\n(\w+)', r'\1\2', text)
        
        # Fix spacing after periods
        text = re.sub(r'\.([A-Z])', r'. \1', text)
        
        # Fix multiple spaces
        text = re.sub(r' +', ' ', text)
        
        return text
    
    def _enhance_academic_text(self, text: str) -> str:
        """Enhance processing for academic content"""
        # Preserve academic terms with hyphens
        text = re.sub(r'([a-z])-([a-z])', r'\1-\2', text)  # Keep valid hyphens
        
        # Handle common academic patterns
        for pattern, replacement in self.academic_patterns:
            text = re.sub(pattern, replacement, text)
        
        return text
    
    def _remove_stopwords(self, text: str) -> str:
        """Remove stopwords from text"""
        words = text.split()
        filtered_words = [word for word in words if word.lower() not in self.stop_words]
        return ' '.join(filtered_words)
    
    def _lemmatize_text(self, text: str) -> str:
        """Lemmatize text"""
        words = text.split()
        lemmatized = [self.lemmatizer.lemmatize(word) for word in words]
        return ' '.join(lemmatized)
    
    def _extract_metadata(self, text: str, doc_type: str) -> Dict[str, Any]:
        """Extract metadata from text content"""
        metadata = {
            'word_count': len(text.split()),
            'char_count': len(text),
            'line_count': len(text.splitlines()),
            'sentence_count': len(nltk.sent_tokenize(text)) if text else 0,
        }
        
        # Extract entities if possible
        if nlp is not None and len(text) < 50000:  # Limit for performance
            metadata['entities'] = self.extract_entities(text)
        
        # Extract key phrases
        metadata['key_phrases'] = self.extract_key_phrases(text, top_n=5)
        
        # Check for academic content indicators
        metadata['has_citations'] = bool(re.search(r'\[\d+\]|\(\w+ et al\.', text))
        metadata['has_equations'] = bool(re.search(r'[\+\-\*\/\^=]|\$.*\$', text))
        
        return metadata
    
    def _generate_document_id(self, file_path: Path, text: str) -> str:
        """Generate unique document ID"""
        content_hash = hashlib.md5(text[:1000].encode()).hexdigest()[:8]
        path_hash = hashlib.md5(str(file_path).encode()).hexdigest()[:8]
        return f"{path_hash}_{content_hash}"
    
    def _load_academic_patterns(self) -> List[tuple]:
        """Load patterns for academic content enhancement"""
        return [
            (r'Fig\.?\s*\d+', 'FIGURE_REFERENCE'),  # Figure references
            (r'Table\s*\d+', 'TABLE_REFERENCE'),     # Table references
            (r'Eq\.?\s*\(\d+\)', 'EQUATION_REFERENCE'),  # Equation references
            (r'Section\s*\d+(\.\d+)*', 'SECTION_REFERENCE'),  # Section references
            (r'Chapter\s*\d+', 'CHAPTER_REFERENCE'),  # Chapter references
            (r'Theorem\s*\d+(\.\d+)*', 'THEOREM_REFERENCE'),  # Theorem references
            (r'Lemma\s*\d+(\.\d+)*', 'LEMMA_REFERENCE'),  # Lemma references
            (r'Proof\.?\s*', 'PROOF_START'),  # Proof indicators
            (r'Q\.E\.D\.?|□|∎', 'PROOF_END'),  # Proof end markers
            (r'Definition\s*\d+(\.\d+)*', 'DEFINITION_REFERENCE'),  # Definitions
            (r'Remark\s*\d+', 'REMARK_REFERENCE'),  # Remarks
            (r'Note\s*\d+', 'NOTE_REFERENCE'),  # Notes
            (r'Example\s*\d+', 'EXAMPLE_REFERENCE'),  # Examples
        ]